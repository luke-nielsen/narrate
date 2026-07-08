"""PowerPoint presentation renderer (requires the ``pptx`` extra)."""

from __future__ import annotations

import io
from typing import ClassVar

from narrate.exceptions import RenderError
from narrate.models import ArtifactKind
from narrate.rendering.base import Renderer, RenderRequest, RenderResult
from narrate.rendering.builtin.decks import parse_deck

__all__ = ["PptxRenderer"]


class PptxRenderer(Renderer):
    """Render slide-deck sources into ``.pptx`` presentations.

    Uses ``python-pptx``; install with ``pip install narrate[pptx]``.
    Consumes the same slide-deck JSON as the HTML slides renderer, so one
    saved artifact yields both outputs.
    """

    name: ClassVar[str] = "pptx"
    input_kinds: ClassVar[frozenset[ArtifactKind]] = frozenset({ArtifactKind.SLIDE_DECK})
    output_kind: ClassVar[ArtifactKind] = ArtifactKind.PRESENTATION
    description: ClassVar[str] = "Render slide decks as PowerPoint (.pptx) presentations"

    def render(self, request: RenderRequest) -> RenderResult:
        """Produce the PowerPoint file.

        Raises:
            RenderError: If ``python-pptx`` is not installed or the deck
                document is malformed.
        """
        self.check_input(request)
        try:
            from pptx import Presentation
        except ImportError as exc:
            msg = "python-pptx is required for the pptx renderer: pip install 'narrate[pptx]'"
            raise RenderError(msg) from exc

        deck = parse_deck(request.json())
        presentation = Presentation()

        title_layout = presentation.slide_layouts[0]
        title_slide = presentation.slides.add_slide(title_layout)
        title_slide.shapes.title.text = deck.title  # pyright: ignore[reportOptionalMemberAccess]
        if deck.subtitle and len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = deck.subtitle  # pyright: ignore[reportAttributeAccessIssue]

        bullet_layout = presentation.slide_layouts[1]
        for slide_spec in deck.slides:
            slide = presentation.slides.add_slide(bullet_layout)
            slide.shapes.title.text = slide_spec.title  # pyright: ignore[reportOptionalMemberAccess]
            body = slide.placeholders[1].text_frame  # pyright: ignore[reportAttributeAccessIssue]
            for position, bullet in enumerate(slide_spec.bullets):
                paragraph = body.paragraphs[0] if position == 0 else body.add_paragraph()
                paragraph.text = bullet
            if slide_spec.notes:
                slide.notes_slide.notes_text_frame.text = slide_spec.notes  # pyright: ignore[reportOptionalMemberAccess]

        buffer = io.BytesIO()
        presentation.save(buffer)
        return RenderResult(
            content=buffer.getvalue(),
            kind=self.output_kind,
            media_type=(
                "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ),
            suffix=".pptx",
            metadata={"slides": len(deck.slides) + 1, "title": deck.title},
        )
